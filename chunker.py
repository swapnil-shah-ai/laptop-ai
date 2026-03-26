"""
Laptop AI - Document Chunker
Extracts text from any supported file type and chunks it intelligently.
Each chunk carries metadata: source file, page/sheet/slide number, position.
"""

import csv
import json
import io
from pathlib import Path
from typing import Generator
from rich.console import Console

console = Console()


def extract_and_chunk(file_info: dict, chunk_size: int = 800, chunk_overlap: int = 100) -> list[dict]:
    """
    Extract text from a file and split into chunks with metadata.
    Returns list of {text, metadata} dicts.
    """
    ext = file_info["extension"]
    path = file_info["path"]

    try:
        if ext == ".pdf":
            return _chunk_pdf(path, chunk_size, chunk_overlap)
        elif ext == ".docx":
            return _chunk_docx(path, chunk_size, chunk_overlap)
        elif ext in (".xlsx", ".xls"):
            return _chunk_excel(path, chunk_size, chunk_overlap)
        elif ext == ".pptx":
            return _chunk_pptx(path, chunk_size, chunk_overlap)
        elif ext == ".csv":
            return _chunk_csv(path, chunk_size, chunk_overlap)
        elif ext == ".json":
            return _chunk_json(path, chunk_size, chunk_overlap)
        elif ext in (".txt", ".md", ".log", ".py", ".js", ".ts", ".html", ".css", ".yaml", ".yml"):
            return _chunk_text(path, chunk_size, chunk_overlap)
        else:
            return []
    except Exception as e:
        # Silently skip files that can't be read
        return []


def _make_chunks(text: str, source: str, chunk_size: int, chunk_overlap: int, extra_meta: dict = None) -> list[dict]:
    """Generic text splitter with overlap."""
    if not text or not text.strip():
        return []

    # Hard cap: truncate any text block to ~6000 chars per chunk max
    # nomic-embed-text has 8192 token limit, this keeps us safely under
    MAX_CHUNK_CHARS = 6000

    # Split by paragraphs first, then merge to chunk_size
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    
    chunks = []
    current_chunk = []
    current_length = 0

    for para in paragraphs:
        para_length = len(para.split())  # approximate token count by words
        
        if current_length + para_length > chunk_size and current_chunk:
            chunk_text = "\n".join(current_chunk)
            if len(chunk_text) > MAX_CHUNK_CHARS:
                chunk_text = chunk_text[:MAX_CHUNK_CHARS]
            meta = {"source": source, "chunk_index": len(chunks)}
            if extra_meta:
                meta.update(extra_meta)
            chunks.append({"text": chunk_text, "metadata": meta})
            
            # Keep overlap - take last few paragraphs
            overlap_text = ""
            overlap_paras = []
            for p in reversed(current_chunk):
                if len(overlap_text.split()) + len(p.split()) <= chunk_overlap:
                    overlap_paras.insert(0, p)
                    overlap_text = "\n".join(overlap_paras)
                else:
                    break
            
            current_chunk = overlap_paras
            current_length = len(overlap_text.split())
        
        current_chunk.append(para)
        current_length += para_length

    # Last chunk
    if current_chunk:
        chunk_text = "\n".join(current_chunk)
        meta = {"source": source, "chunk_index": len(chunks)}
        if extra_meta:
            meta.update(extra_meta)
        chunks.append({"text": chunk_text, "metadata": meta})

    return chunks


# ---- File type handlers ----

def _chunk_pdf(path: str, chunk_size: int, chunk_overlap: int) -> list[dict]:
    from pypdf import PdfReader
    reader = PdfReader(path)
    all_chunks = []
    
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            chunks = _make_chunks(
                text, path, chunk_size, chunk_overlap,
                extra_meta={"page": i + 1, "file_type": "pdf"}
            )
            all_chunks.extend(chunks)
    
    return all_chunks


def _chunk_docx(path: str, chunk_size: int, chunk_overlap: int) -> list[dict]:
    from docx import Document
    doc = Document(path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    
    # Also grab tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text += "\n" + row_text
    
    return _make_chunks(text, path, chunk_size, chunk_overlap, {"file_type": "docx"})


def _chunk_excel(path: str, chunk_size: int, chunk_overlap: int) -> list[dict]:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    all_chunks = []
    
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        headers = None
        
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            # Convert None to empty string
            row_values = [str(v) if v is not None else "" for v in row]
            
            # Skip completely empty rows
            if not any(v.strip() for v in row_values):
                continue
            
            if i == 0:
                headers = row_values
                continue
            
            if headers:
                # Create key-value representation for each row
                row_text = " | ".join(
                    f"{headers[j]}: {row_values[j]}"
                    for j in range(min(len(headers), len(row_values)))
                    if row_values[j].strip()
                )
            else:
                row_text = " | ".join(v for v in row_values if v.strip())
            
            if row_text.strip():
                rows.append(row_text)
        
        # Group rows into chunks
        if rows:
            header_line = f"[Sheet: {sheet_name}]"
            if headers:
                header_line += f"\nColumns: {' | '.join(h for h in headers if h.strip())}"
            
            text = header_line + "\n" + "\n".join(rows)
            chunks = _make_chunks(
                text, path, chunk_size, chunk_overlap,
                extra_meta={"sheet": sheet_name, "file_type": "excel"}
            )
            all_chunks.extend(chunks)
    
    wb.close()
    return all_chunks


def _chunk_pptx(path: str, chunk_size: int, chunk_overlap: int) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(path)
    all_chunks = []
    
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
        
        if texts:
            slide_text = f"[Slide {i + 1}]\n" + "\n".join(texts)
            chunks = _make_chunks(
                slide_text, path, chunk_size, chunk_overlap,
                extra_meta={"slide": i + 1, "file_type": "pptx"}
            )
            all_chunks.extend(chunks)
    
    return all_chunks


def _chunk_csv(path: str, chunk_size: int, chunk_overlap: int) -> list[dict]:
    rows = []
    headers = None
    
    # Try different encodings
    for encoding in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
        try:
            with open(path, "r", encoding=encoding) as f:
                reader = csv.reader(f)
                for i, row in enumerate(reader):
                    if i == 0:
                        headers = row
                        continue
                    if headers:
                        row_text = " | ".join(
                            f"{headers[j]}: {row[j]}"
                            for j in range(min(len(headers), len(row)))
                            if row[j].strip()
                        )
                    else:
                        row_text = " | ".join(v for v in row if v.strip())
                    if row_text.strip():
                        rows.append(row_text)
            break
        except (UnicodeDecodeError, csv.Error):
            continue

    if rows:
        text = "\n".join(rows)
        return _make_chunks(text, path, chunk_size, chunk_overlap, {"file_type": "csv"})
    return []


def _chunk_json(path: str, chunk_size: int, chunk_overlap: int) -> list[dict]:
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        text = json.dumps(data, indent=2, default=str)
        return _make_chunks(text, path, chunk_size, chunk_overlap, {"file_type": "json"})
    except (json.JSONDecodeError, UnicodeDecodeError):
        return []


def _chunk_text(path: str, chunk_size: int, chunk_overlap: int) -> list[dict]:
    for encoding in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
        try:
            with open(path, "r", encoding=encoding) as f:
                text = f.read()
            ext = Path(path).suffix.lower()
            return _make_chunks(text, path, chunk_size, chunk_overlap, {"file_type": ext.lstrip(".")})
        except UnicodeDecodeError:
            continue
    return []
